from pptx import Presentation
from pptx.util import Inches

# 新しいプレゼンテーションを作成
prs = Presentation()

# スライドマスターの編集
# title_slide_layout = prs.slide_master.slide_layouts[0]
# placeholder = title_slide_layout.placeholders[1]
# placeholder.left = Inches(2.5)
# placeholder.top = Inches(4)


# タイトルスライドを追加
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "プレゼンテーションのタイトル"
subtitle.text = "サブタイトル"

prs.save("new_presentation.pptx")
